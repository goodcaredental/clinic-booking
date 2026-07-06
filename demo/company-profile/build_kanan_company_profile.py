"""Generate the Kanan Digital Enterprise company profile.

Outputs ./demo/company-profile/Kanan_Company_Profile.pptx — 10 slides,
Kanan brand styling (navy + gold + warm-white per the brand guidelines).
Content mirrors ./kanan-company-profile.md — edit the .md for copy, then re-run.

Different scope from the sales deck: this is about *who Kanan is*, not what
one product does. Prospects, partners, potential clients evaluating whether
to work with us at all should be able to read this and understand our
positioning + engagement model without seeing a product screenshot.

Run:  python3 demo/company-profile/build_kanan_company_profile.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── Kanan brand colors (must match brand guidelines memory) ─────────────
NAVY        = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_LIGHT  = RGBColor(0x2E, 0x43, 0x74)
GOLD        = RGBColor(0xC9, 0xA2, 0x27)
GOLD_LIGHT  = RGBColor(0xE3, 0xC7, 0x6A)
WARM_WHITE  = RGBColor(0xF4, 0xF1, 0xEA)
GREY        = RGBColor(0x6B, 0x72, 0x80)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

HEAD_FONT = "Helvetica"
BODY_FONT = "Helvetica"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)


# ─── Helpers (same primitives as build_kanan_sales_deck.py) ─────────────

def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(shape, color)
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, italic=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=BODY_FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, *,
                size=14, color=NAVY, bullet_color=GOLD,
                line_gap_pt=6, bullet_char="●"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(line_gap_pt)
        b = p.add_run()
        b.text = f"{bullet_char}  "
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        b.font.name = BODY_FONT
        t = p.add_run()
        t.text = item
        t.font.size = Pt(size)
        t.font.color.rgb = color
        t.font.name = BODY_FONT
    return box


def add_speaker_note(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_footer(slide, page_num, total):
    add_text(slide, MARGIN, SLIDE_H - Inches(0.4),
             Inches(6), Inches(0.3),
             "Kanan Digital Enterprise · kanan.my",
             size=10, color=GREY, font=BODY_FONT)
    add_text(slide, SLIDE_W - MARGIN - Inches(1), SLIDE_H - Inches(0.4),
             Inches(1), Inches(0.3), f"{page_num} / {total}",
             size=10, color=GREY, align=PP_ALIGN.RIGHT, font=BODY_FONT)


def add_screenshot_placeholder(slide, left, top, width, height, caption=""):
    box = add_rect(slide, left, top, width, height, WARM_WHITE, line_color=GOLD_LIGHT)
    box.line.width = Pt(1.5)
    add_text(slide, left, top, width, height, "[ screenshot ]",
             size=11, italic=True, color=GREY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=BODY_FONT)
    if caption:
        add_text(slide, left, top + height + Inches(0.05), width, Inches(0.3),
                 caption, size=10, italic=True, color=GREY,
                 align=PP_ALIGN.CENTER, font=BODY_FONT)


def slide_title_only(prs, page, total, title, lead=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WARM_WHITE)
    add_rect(s, MARGIN, Inches(0.55), Inches(0.5), Inches(0.05), GOLD)
    add_text(s, MARGIN, Inches(0.7), SLIDE_W - 2 * MARGIN, Inches(0.8),
             title, size=32, bold=True, color=NAVY, font=HEAD_FONT)
    if lead:
        add_text(s, MARGIN, Inches(1.6), SLIDE_W - 2 * MARGIN, Inches(0.6),
                 lead, size=16, italic=True, color=NAVY_LIGHT)
    add_footer(s, page, total)
    return s


# ─── Slide builders ───────────────────────────────────────────────────────

def s1_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WARM_WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.15), GOLD)
    # Wordmark
    add_text(s, 0, Inches(2.2), SLIDE_W, Inches(0.6),
             "kanan", size=80, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, font="Georgia")
    add_text(s, 0, Inches(3.2), SLIDE_W, Inches(0.5),
             "your trusted right hand",
             size=18, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.2), SLIDE_W, Inches(0.6),
             "Kanan Digital Enterprise",
             size=22, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.9), SLIDE_W, Inches(0.4),
             "Company profile · 2026",
             size=14, color=GREY, italic=True, align=PP_ALIGN.CENTER)
    add_rect(s, 0, SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), NAVY)
    add_text(s, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.3),
             "kanan.my · hello@kanan.my", size=11, color=WARM_WHITE,
             align=PP_ALIGN.CENTER)


def s2_who_we_are(prs, total):
    s = slide_title_only(prs, 2, total,
                         "A Malaysian software studio, built by operators.")
    # Etymology block — italic, indented
    add_text(s, MARGIN, Inches(1.9), SLIDE_W - 2 * MARGIN, Inches(1.0),
             "Kanan — 'KAH-nahn' — means right in Malay. As in right hand. "
             "Dependable. Close by. The hand that carries what the main one can't.",
             size=15, italic=True, color=NAVY_LIGHT)
    add_bullets(s, MARGIN, Inches(3.3), SLIDE_W - 2 * MARGIN, Inches(3.5), [
        "Founded in Malaysia by [Founder name 1] and [Founder name 2].",
        "We build software agents — small, focused tools that automate the parts of a business that shouldn't need a human anymore.",
        "Every product we ship is one we'd use ourselves in a business we could run.",
        "Not a consultancy chasing hours. Not a startup chasing venture rounds. An operating studio building products that fund themselves.",
    ], size=15, line_gap_pt=12)
    add_speaker_note(s, "Personalize. Founder names + real backgrounds if possible. "
                        "Own the 'operator, not consultant' line — this is where prospects decide "
                        "whether we're the right kind of vendor.")


def s3_what_we_build(prs, total):
    s = slide_title_only(prs, 3, total,
                         "AI agents for Malaysian businesses.",
                         "One product per business problem. Deployed, supported, owned end-to-end.")
    # 5-card icon strip
    products = [
        ("🩺", "Clinic Booking",   "Patient self-booking + nurse queue + owner dashboard.",   "shipped"),
        ("📅", "Personal Assistant", "Meeting coordination + follow-ups + summary emails.",  "building"),
        ("📋", "Purchase Order",   "PO extraction, LMW/FTZ checks, accounting sync.",         "building"),
        ("📄", "PDF Translation",  "Long-form document translation with layout kept intact.", "building"),
        ("📢", "Marketing Agent",  "Content pipeline + audit trail for regulated industries.","building"),
    ]
    card_top = Inches(2.6)
    total_cards = len(products)
    total_gap = Inches(0.15) * (total_cards - 1)
    card_w = (SLIDE_W - 2 * MARGIN - total_gap) // total_cards
    card_h = Inches(3.0)
    for i, (icon, name, desc, status) in enumerate(products):
        x = MARGIN + i * (card_w + Inches(0.15))
        add_rect(s, x, card_top, card_w, card_h, WHITE, line_color=NAVY_LIGHT)
        status_color = GOLD if status == "shipped" else GREY
        # Top accent bar
        add_rect(s, x, card_top, card_w, Inches(0.1), status_color)
        # Icon
        add_text(s, x, card_top + Inches(0.35), card_w, Inches(0.7),
                 icon, size=36, color=NAVY, align=PP_ALIGN.CENTER)
        # Product name
        add_text(s, x + Inches(0.15), card_top + Inches(1.2), card_w - Inches(0.3), Inches(0.5),
                 name, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # Status pill
        add_text(s, x + Inches(0.15), card_top + Inches(1.65), card_w - Inches(0.3), Inches(0.3),
                 status.upper(), size=9, bold=True, color=status_color,
                 align=PP_ALIGN.CENTER)
        # Description
        add_text(s, x + Inches(0.2), card_top + Inches(2.05), card_w - Inches(0.4), Inches(0.9),
                 desc, size=10, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(6.05), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Every product starts with a real client's real workflow. We don't build for hypothetical markets.",
             size=12, italic=True, color=NAVY_LIGHT, align=PP_ALIGN.CENTER)


def s4_approach(prs, total):
    s = slide_title_only(prs, 4, total,
                         "Small team. Real users. No middlemen.")
    blocks = [
        ("We build in weeks, not months.",
         "Prototype on your real data inside two weeks. If it doesn't earn the demo, we stop before you're committed."),
        ("We stay after launch.",
         "No 3-month handover then radio silence. We're on WhatsApp when something breaks — and when you have a question."),
        ("Your data stays yours.",
         "Hosted in Malaysia where possible (Supabase + Vercel). We don't train external AI models on your operational data. PDPA-aware from day one."),
        ("We don't force upgrades.",
         "Feature tiers exist. Seat caps exist. But we top up seats on request — no forced tier bump just to add one more nurse."),
    ]
    grid_top = Inches(2.2)
    cell_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) // 2
    cell_h = Inches(2.2)
    for i, (title, body) in enumerate(blocks):
        row, col = divmod(i, 2)
        x = MARGIN + col * (cell_w + Inches(0.3))
        y = grid_top + row * (cell_h + Inches(0.3))
        add_rect(s, x, y, cell_w, cell_h, WHITE, line_color=NAVY_LIGHT)
        add_rect(s, x + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.04), GOLD)
        add_text(s, x + Inches(0.3), y + Inches(0.5), cell_w - Inches(0.6), Inches(0.5),
                 title, size=15, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), y + Inches(1.05), cell_w - Inches(0.6), Inches(1.1),
                 body, size=12, color=GREY)


def s5_case_study(prs, total):
    s = slide_title_only(prs, 5, total,
                         "Kanan Clinic Booking — a case study.",
                         "Our first shipped product. Live at demo clinics in KL, more onboarding through 2026.")
    add_bullets(s, MARGIN, Inches(2.4), Inches(7.5), Inches(4.5), [
        "Problem: Nurses at small dental clinics spent 30–50 min/day on WhatsApp — confirming, reminding, chasing. Owners had zero visibility.",
        "What we built: Multi-tenant SaaS with patient self-booking (EN / 中文 / BM), nurse queue with pre-written WhatsApp templates, terminal-kiosk with PIN-per-action, owner dashboard with full audit trail.",
        "Two tiers — Standard for single-doctor clinics; Premium adds a room flow (nurse check-in → doctor check-out with treatment logged), performance analytics, chair utilization heatmap.",
        "Outcome (indicative, ask for latest): no-show rate down, nurse WhatsApp time cut ~60%, owners can name their busiest day of the week for the first time.",
    ], size=13, line_gap_pt=10)
    # 2-panel screenshot area on the right
    add_screenshot_placeholder(s, Inches(8.5), Inches(2.4), Inches(4.2), Inches(2.0),
                               "Owner dashboard")
    add_screenshot_placeholder(s, Inches(8.5), Inches(4.9), Inches(4.2), Inches(2.0),
                               "Terminal lockscreen")
    add_speaker_note(s, "Same code runs both demo URLs — the tier flag is the difference. That's the point.")


def s6_how_we_work(prs, total):
    s = slide_title_only(prs, 6, total, "Start to finish — no surprises")
    steps = [
        ("1", "Talk to us", "WhatsApp or a short call. Tell us what's slowing you down."),
        ("2", "Scoping session", "Longer conversation. No sales deck."),
        ("3", "Workflow map", "Map your current ops + where software fits."),
        ("4", "Written proposal", "1-page scope, timeline, price. Push back where it doesn't fit."),
        ("5", "Working demo", "On your real data, in weeks not months."),
        ("6", "Commercial agreed", "Full commercial agreed once the demo proves out."),
        ("7", "Setup + training", "Team trained, real users on the system. WhatsApp throughout."),
        ("8", "Ongoing support", "We're still here. That's the whole business model."),
    ]
    grid_top = Inches(2.0)
    cell_w = (SLIDE_W - 2 * MARGIN - Inches(0.3) * 3) // 4
    cell_h = Inches(2.1)
    for i, (n, title, body) in enumerate(steps):
        row, col = divmod(i, 4)
        x = MARGIN + col * (cell_w + Inches(0.3))
        y = grid_top + row * (cell_h + Inches(0.3))
        add_rect(s, x, y, cell_w, cell_h, WHITE, line_color=NAVY_LIGHT)
        add_rect(s, x + Inches(0.25), y + Inches(0.25), Inches(0.5), Inches(0.04), GOLD)
        add_text(s, x + Inches(0.25), y + Inches(0.4), cell_w - Inches(0.5), Inches(0.5),
                 n, size=28, bold=True, color=NAVY, font="Georgia")
        add_text(s, x + Inches(0.25), y + Inches(0.95), cell_w - Inches(0.5), Inches(0.45),
                 title, size=13, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), y + Inches(1.4), cell_w - Inches(0.5), Inches(0.7),
                 body, size=10, color=GREY)
    add_speaker_note(s, "Lead with step 8. 'We don't ghost you after launch' is the differentiator.")


def s7_team(prs, total):
    s = slide_title_only(prs, 7, total,
                         "Two founders. That's the whole team, for now.")
    # Two founder cards
    card_w = Inches(5.5)
    card_h = Inches(2.6)
    card_top = Inches(2.4)
    gap = Inches(0.3)
    left_start = (SLIDE_W - (card_w * 2 + gap)) // 2
    founders = [
        ("[Founder name 1]", "Co-founder · [Product & Engineering]",
         "[1-line bio — background, focus area, what they own inside Kanan]"),
        ("[Founder name 2]", "Co-founder · [Client & Ops]",
         "[1-line bio — background, focus area, what they own inside Kanan]"),
    ]
    for i, (name, role, bio) in enumerate(founders):
        x = left_start + i * (card_w + gap)
        add_rect(s, x, card_top, card_w, card_h, WHITE, line_color=NAVY_LIGHT)
        add_rect(s, x, card_top, card_w, Inches(0.15), GOLD)
        add_text(s, x + Inches(0.4), card_top + Inches(0.4), card_w - Inches(0.8), Inches(0.6),
                 name, size=20, bold=True, color=NAVY)
        add_text(s, x + Inches(0.4), card_top + Inches(1.05), card_w - Inches(0.8), Inches(0.4),
                 role, size=13, italic=True, color=GOLD)
        add_text(s, x + Inches(0.4), card_top + Inches(1.55), card_w - Inches(0.8), Inches(0.9),
                 bio, size=12, color=GREY)
    add_text(s, MARGIN, Inches(5.4), SLIDE_W - 2 * MARGIN, Inches(0.5),
             "We answer our own WhatsApp. If you sign with us, you won't be handed off to a junior.",
             size=14, italic=True, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(6.05), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Hiring plan: one engineer + one ops person once we've supported five paying customers through a full cycle. Slow by design.",
             size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)


def s8_trust(prs, total):
    s = slide_title_only(prs, 8, total, "Your data stays yours",
                         "Trust before technology.")
    add_bullets(s, MARGIN, Inches(2.4), SLIDE_W - 2 * MARGIN, Inches(4), [
        "Hosted in Malaysia wherever possible (Supabase + Vercel). Data stays within the ASEAN legal envelope.",
        "PDPA-aware — patient / customer data export any time, schema visible on request, no third-party analytics.",
        "We don't train external AI models on your operational data. Inference runs against your data only during the request.",
        "You own your backups. Owners can download a full CSV any day; daily auto-email backup available.",
        "Clear scope line — clinical records / X-rays / sterilization logs stay in your EMR, not in our system.",
    ], size=14, line_gap_pt=10)
    add_text(s, MARGIN, Inches(6.5), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Kanan Digital Enterprise · SSM Enterprise registration in progress · Kuala Lumpur, Malaysia",
             size=11, italic=True, color=GREY)


def s9_wont_do(prs, total):
    s = slide_title_only(prs, 9, total, "Scope discipline is a feature.",
                         "We say no to things that would dilute the products for the businesses that actually use them.")
    add_bullets(s, MARGIN, Inches(2.6), SLIDE_W - 2 * MARGIN, Inches(4), [
        "We don't build medical record systems. Clinic Booking is an ops layer, not an EMR — if you need an EMR we'll point you at one that does it well.",
        "We don't do fixed-price 'digital transformation' projects. We ship one working thing, then the next.",
        "We don't sell AI as a feature. We ship products that happen to use AI where it earns its keep. If a simple form works better than an agent, we build the form.",
        "We don't take on projects we can't support for two years. Small team, deep bench per product. When we take you on, we're on for the whole ride.",
    ], size=15, line_gap_pt=12, bullet_char="—")


def s10_contact(prs, total):
    s = slide_title_only(prs, 10, total, "Let's have a conversation.")
    # Contact block
    band_top = Inches(2.2)
    add_rect(s, MARGIN, band_top, SLIDE_W - 2 * MARGIN, Inches(2.0), NAVY)
    add_text(s, MARGIN, band_top + Inches(0.4), SLIDE_W - 2 * MARGIN, Inches(0.6),
             "WhatsApp +60 12-347 8126",
             size=22, bold=True, color=WARM_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, band_top + Inches(1.05), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "hello@kanan.my   ·   kanan.my",
             size=16, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

    # Demo URLs
    demo_top = Inches(4.6)
    add_text(s, MARGIN, demo_top, SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Product demos",
             size=11, bold=True, color=GREY, align=PP_ALIGN.CENTER,
             font=BODY_FONT)
    card_w = Inches(5.5)
    card_h = Inches(1.4)
    card_top = demo_top + Inches(0.5)
    gap = Inches(0.3)
    left_start = (SLIDE_W - (card_w * 2 + gap)) // 2
    for i, (title, url) in enumerate([
        ("Clinic Booking · Standard", "standard-demo.kanan.my"),
        ("Clinic Booking · Premium", "premium-demo.kanan.my"),
    ]):
        x = left_start + i * (card_w + gap)
        accent = NAVY_LIGHT if i == 0 else GOLD
        add_rect(s, x, card_top, card_w, card_h, WHITE, line_color=accent)
        add_rect(s, x, card_top, card_w, Inches(0.12), accent)
        add_text(s, x, card_top + Inches(0.3), card_w, Inches(0.4),
                 title, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, x, card_top + Inches(0.75), card_w, Inches(0.4),
                 url, size=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(6.9), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "No pitches on the first call. Just tell us what's not working, and we'll tell you honestly whether software is the right answer.",
             size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# ─── Build it ─────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        s1_cover,
        s2_who_we_are,
        s3_what_we_build,
        s4_approach,
        s5_case_study,
        s6_how_we_work,
        s7_team,
        s8_trust,
        s9_wont_do,
        s10_contact,
    ]
    total = len(builders)
    builders[0](prs)
    for i, fn in enumerate(builders[1:], start=2):
        fn(prs, total)

    out = Path(__file__).parent / "Kanan_Company_Profile.pptx"
    prs.save(out)
    print(f"Wrote {out} ({total} slides)")


if __name__ == "__main__":
    build()
